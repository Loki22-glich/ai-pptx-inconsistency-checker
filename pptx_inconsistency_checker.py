import os
import json
from pptx import Presentation
from PIL import Image
import pytesseract
import google.generativeai as genai

# ---------------- CONFIG ----------------
API_KEY = os.getenv("GEMINI_API_KEY")
if not API_KEY:
    raise ValueError("Please set GEMINI_API_KEY as an environment variable.")
genai.configure(api_key=API_KEY)
MODEL_NAME = "gemini-2.5-flash"
# -----------------------------------------

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slide_texts = {}
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slide_texts[i] = "\n".join(texts)
    return slide_texts

def extract_text_from_images(image_paths):
    ocr_results = {}
    for path in image_paths:
        try:
            img = Image.open(path)
            text = pytesseract.image_to_string(img)
            if text.strip():
                ocr_results[os.path.basename(path)] = text
        except Exception as e:
            print(f"OCR failed for {path}: {e}")
    return ocr_results

def build_facts(slide_texts, ocr_texts):
    facts = []
    for slide_no, text in slide_texts.items():
        for line in text.split("\n"):
            if line.strip():
                facts.append({"slide": slide_no, "text": line.strip()})
    for img_name, text in ocr_texts.items():
        for line in text.split("\n"):
            if line.strip():
                facts.append({"slide": img_name, "text": line.strip()})
    return facts

def check_inconsistencies_with_gemini(facts):
    model = genai.GenerativeModel(MODEL_NAME)
    prompt = f"""
You are an AI that finds inconsistencies in PowerPoint presentations.

Facts (JSON):
{json.dumps(facts, indent=2)}

Task:
1. Identify factual or logical contradictions across these statements.
2. Contradictions may be:
   - Numeric (different numbers for same metric)
   - Textual (claims that can't both be true)
   - Timeline/date mismatches
3. For each inconsistency, return:
   - "slides": list of slide numbers or image names
   - "type": "numeric", "textual", or "timeline"
   - "description": short explanation

Output JSON only.
"""
    response = model.generate_content(prompt)
    raw = response.text.strip()
    if raw.startswith("```"):
        raw = raw.strip("`").lstrip("json").strip()
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        print("Gemini output not valid JSON. Raw output:\n", raw)
        return []

def find_images_in_same_folder(pptx_path):
    folder = os.path.dirname(pptx_path)
    image_exts = {".jpg", ".jpeg", ".png"}
    images = []
    for file in os.listdir(folder):
        if os.path.splitext(file)[1].lower() in image_exts:
            images.append(os.path.join(folder, file))
    return images

def main():
    import sys
    if len(sys.argv) < 2:
        print("Usage: python pptx_inconsistency_checker.py <pptx_file>")
        return

    pptx_file = sys.argv[1]
    if not os.path.exists(pptx_file):
        print(f"Error: {pptx_file} not found")
        return

    print("Extracting text from PPTX...")
    slide_texts = extract_text_from_pptx(pptx_file)
    print(f"Extracted {len(slide_texts)} slides.")

    print("Searching for images in same folder...")
    image_files = find_images_in_same_folder(pptx_file)
    print(f"Found {len(image_files)} images for OCR.")

    ocr_texts = extract_text_from_images(image_files)
    print(f"OCR extracted text from {len(ocr_texts)} images.")

    facts = build_facts(slide_texts, ocr_texts)
    print(f"Built {len(facts)} facts.")

    print("Checking inconsistencies with Gemini...")
    issues = check_inconsistencies_with_gemini(facts)

    # Save results
    with open("inconsistencies.json", "w") as f:
        json.dump(issues, f, indent=2)

    # Pretty print
    print("\n=== Inconsistencies Found ===")
    if not issues:
        print("No inconsistencies detected.")
    else:
        for idx, issue in enumerate(issues, start=1):
            slides = ", ".join(str(s) for s in issue.get("slides", []))
            print(f"{idx}. [{issue.get('type', 'unknown').upper()}] Slides: {slides} → {issue.get('description')}")

    print("\nFull results saved to inconsistencies.json")

if __name__ == "__main__":
    main()
